from __future__ import annotations

import hashlib
import json
from io import BytesIO

import pytest
from pptx import Presentation
from pypdf import PdfReader

from app.services.training_presentation_renderer import (
    PDF_CONTENT_TYPE,
    PPTX_CONTENT_TYPE,
    RENDERER_VERSION,
    render_presentation,
    verify_manifest,
)


def _manifest() -> dict:
    value = {
        "manifest_version": "nace-training-presentation-manifest-v1",
        "contract_version": "nace-training-presentation-contract-v1",
        "contract_hash": "b" * 64,
        "template_version": "osgb-training-presentation-template-v1",
        "output_formats": ["pptx", "pdf"],
        "training": {
            "training_id": 101,
            "company_id": 35,
            "title": "Türkçe NACE Uyumlu İSG Eğitimi",
        },
        "source_registry": [
            {
                "source_id": "tr-law-6331",
                "title": "6331 sayılı İş Sağlığı ve Güvenliği Kanunu",
            }
        ],
        "slides": [
            {
                "position": 1,
                "section_id": "cover",
                "title": "Türkçe NACE Uyumlu İSG Eğitimi",
                "source_refs": ["tr-law-6331"],
                "content_blocks": [
                    {
                        "type": "nace_identity",
                        "nace_code": "62.01.01",
                        "nace_description": "Bilgisayar programlama faaliyetleri",
                    },
                    {"type": "hazard_class", "value": "Az Tehlikeli"},
                    {"type": "company_logo_placeholder", "value": None},
                ],
                "speaker_notes_required": True,
                "approval_required": True,
            },
            {
                "position": 2,
                "section_id": "technical_risks",
                "title": "Teknik risk profili",
                "source_refs": ["tr-law-6331", "controlled_risk_catalog:test"],
                "content_blocks": [
                    {"type": "technical_risk_tag", "value": "display_screen"},
                    {"type": "technical_risk_tag", "value": "server_room"},
                    {"type": "technical_risk_tag", "value": "psychosocial"},
                ],
                "speaker_notes_required": True,
                "approval_required": False,
            },
            {
                "position": 3,
                "section_id": "assessment",
                "title": "Bilgi kontrolü ve değerlendirme",
                "source_refs": ["approved_question_bank:test"],
                "content_blocks": [
                    {"type": "exam_distribution", "foundation": 5, "work_specific": 15},
                    {"type": "exam_workflow_unchanged", "value": True},
                ],
                "speaker_notes_required": True,
                "approval_required": False,
            },
        ],
        "slide_count": 3,
        "approval": {
            "status": "specialist_review_required",
            "required_slide_positions": [1],
        },
        "core_training_unaffected": True,
    }
    value["content_hash"] = hashlib.sha256(
        json.dumps(
            value,
            ensure_ascii=False,
            sort_keys=True,
            separators=(",", ":"),
        ).encode("utf-8")
    ).hexdigest()
    return value


def test_renderer_outputs_are_deterministic_valid_and_reopenable():
    manifest = _manifest()
    first = render_presentation(manifest)
    second = render_presentation(manifest)

    assert first.renderer_version == RENDERER_VERSION
    assert first.slide_count == 3
    assert first.pptx_bytes == second.pptx_bytes
    assert first.pdf_bytes == second.pdf_bytes
    assert first.pptx_bytes.startswith(b"PK")
    assert first.pdf_bytes.startswith(b"%PDF")

    pptx = Presentation(BytesIO(first.pptx_bytes))
    assert len(pptx.slides) == 3
    assert "Türkçe NACE" in pptx.core_properties.title
    for index, slide in enumerate(pptx.slides, start=1):
        notes = slide.notes_slide.notes_text_frame.text
        assert f"Slayt {index}/3" in notes
        assert "değişmez manifestten" in notes

    pdf = PdfReader(BytesIO(first.pdf_bytes))
    assert len(pdf.pages) == 3
    extracted = " ".join(page.extract_text() or "" for page in pdf.pages)
    assert "Türkçe NACE Uyumlu" in extracted
    assert "Teknik risk profili" in extracted
    assert "Mevcut sınav" in extracted


def test_renderer_uses_expected_mime_types():
    assert PPTX_CONTENT_TYPE.endswith("presentationml.presentation")
    assert PDF_CONTENT_TYPE == "application/pdf"


def test_manifest_tampering_fails_before_file_generation():
    manifest = _manifest()
    manifest["slides"][1]["title"] = "Sonradan değiştirildi"
    with pytest.raises(ValueError, match="hash doğrulaması"):
        verify_manifest(manifest)
    with pytest.raises(ValueError, match="hash doğrulaması"):
        render_presentation(manifest)


def test_manifest_requires_contiguous_positions_and_matching_count():
    manifest = _manifest()
    manifest["slides"][2]["position"] = 4
    unsigned = dict(manifest)
    unsigned.pop("content_hash")
    manifest["content_hash"] = hashlib.sha256(
        json.dumps(unsigned, ensure_ascii=False, sort_keys=True, separators=(",", ":")).encode()
    ).hexdigest()
    with pytest.raises(ValueError, match="slayt sırası"):
        verify_manifest(manifest)
