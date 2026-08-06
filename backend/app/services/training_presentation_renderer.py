"""Deterministic PPTX and PDF renderer for frozen NACE presentation manifests."""
from __future__ import annotations

import hashlib
import json
import zipfile
from dataclasses import dataclass
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pypdf import PdfReader
from reportlab.lib.units import inch
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas


SLIDE_WIDTH_IN = 13.333333
SLIDE_HEIGHT_IN = 7.5
PPTX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
PDF_CONTENT_TYPE = "application/pdf"
RENDERER_VERSION = "nace-training-presentation-renderer-v1"

_TEAL = RGBColor(8, 127, 120)
_TEAL_DARK = RGBColor(6, 84, 82)
_MINT = RGBColor(224, 245, 242)
_NAVY = RGBColor(17, 38, 53)
_SLATE = RGBColor(82, 101, 120)
_WHITE = RGBColor(255, 255, 255)
_AMBER = RGBColor(180, 100, 0)

_RISK_LABELS = {
    "display_screen": "Ekranlı araçlarla çalışma",
    "server_room": "Sistem odası elektrik, sıcaklık ve yangın riski",
    "psychosocial": "İş yükü ve psikososyal riskler",
    "working_at_height": "Yüksekte çalışma",
    "excavation": "Kazı çalışmaları",
    "lifting": "Kaldırma ve taşıma işleri",
    "temporary_electricity": "Geçici elektrik tesisatı",
    "cosmetic_chemicals": "Kozmetik kimyasallar",
    "skin_exposure": "Cilt maruziyeti",
    "sterilization": "Sterilizasyon ve hijyen",
    "flammable_gases": "Yanıcı gazlar",
    "atex": "Patlayıcı ortam",
    "confined_space": "Kapalı alan",
    "animal_attack": "Hayvan saldırısı",
    "zoonoses": "Zoonotik hastalıklar",
    "anesthetic_gases": "Anestezik gazlar",
}


@dataclass(frozen=True)
class RenderedPresentation:
    pptx_bytes: bytes
    pdf_bytes: bytes
    slide_count: int
    renderer_version: str = RENDERER_VERSION


def canonical_json_bytes(value: object) -> bytes:
    return json.dumps(
        value,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
    ).encode("utf-8")


def sha256_hex(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def verify_manifest(manifest: dict[str, Any]) -> None:
    if not isinstance(manifest, dict):
        raise ValueError("Sunum manifesti nesne olmalıdır.")
    expected = str(manifest.get("content_hash") or "")
    if len(expected) != 64:
        raise ValueError("Sunum manifest content_hash alanı geçersiz.")
    unsigned = dict(manifest)
    unsigned.pop("content_hash", None)
    actual = sha256_hex(canonical_json_bytes(unsigned))
    if actual != expected:
        raise ValueError("Sunum manifest hash doğrulaması başarısız.")
    slides = manifest.get("slides")
    if not isinstance(slides, list) or not slides:
        raise ValueError("Sunum manifestinde slayt bulunmuyor.")
    positions = [int(item.get("position") or 0) for item in slides if isinstance(item, dict)]
    if positions != list(range(1, len(slides) + 1)):
        raise ValueError("Sunum slayt sırası geçersiz.")
    if int(manifest.get("slide_count") or 0) != len(slides):
        raise ValueError("Sunum slayt sayısı manifest ile uyuşmuyor.")


def _clean(value: object, *, limit: int = 260) -> str:
    text = " ".join(str(value or "").replace("\n", " ").split())
    if len(text) <= limit:
        return text
    return text[: limit - 1].rstrip() + "…"


def _risk_label(value: object) -> str:
    raw = _clean(value, limit=120)
    return _RISK_LABELS.get(raw, raw.replace("_", " ").title())


def _source_map(manifest: dict[str, Any]) -> dict[str, str]:
    result: dict[str, str] = {}
    for item in manifest.get("source_registry") or []:
        if not isinstance(item, dict):
            continue
        source_id = str(item.get("source_id") or "").strip()
        if source_id:
            result[source_id] = _clean(item.get("title") or source_id, limit=120)
    return result


def _source_titles(manifest: dict[str, Any], refs: list[object]) -> list[str]:
    registry = _source_map(manifest)
    titles: list[str] = []
    for ref in refs:
        key = str(ref or "").strip()
        if not key:
            continue
        if key in registry:
            title = registry[key]
        elif key.startswith("training_nace_snapshot:"):
            title = "Dondurulmuş NACE sınıflandırma snapshot'ı"
        elif key.startswith("controlled_training_topics:"):
            title = "Kontrollü beş eğitim konusu"
        elif key.startswith("controlled_risk_catalog:"):
            title = "Kontrollü teknik risk kataloğu"
        elif key.startswith("approved_question_bank:"):
            title = "Onaylı 5 temel + 15 işe özgü sınav içeriği"
        else:
            title = _clean(key, limit=120)
        if title not in titles:
            titles.append(title)
    return titles


def _block_bullets(block: dict[str, Any]) -> list[str]:
    kind = str(block.get("type") or "").strip()
    if kind == "nace_identity":
        return [
            f"NACE: {_clean(block.get('nace_code'))}",
            _clean(block.get("nace_description")),
        ]
    if kind == "hazard_class":
        return [f"Tehlike sınıfı: {_clean(block.get('value'))}"]
    if kind == "training_date":
        return [f"Eğitim tarihi: {_clean(block.get('value')) or 'Kayıtta belirtilmemiş'}"]
    if kind == "company_logo_placeholder":
        return ["Kurum logosu: uzman/kurum onayı bekleniyor"]
    if kind == "topic_objective":
        topic = _clean(block.get("topic"))
        return [f"Öğrenme hedefi: {topic}"] if topic else []
    if kind == "official_source":
        return [f"Resmî kaynak: {_clean(block.get('source_id'))}"]
    if kind == "nace_snapshot":
        return [
            f"NACE kodu: {_clean(block.get('nace_code'))}",
            f"Faaliyet: {_clean(block.get('nace_description'))}",
            f"Tehlike sınıfı: {_clean(block.get('hazard_class'))}",
            f"İçerik profili: {_clean(block.get('content_profile_code'))}",
        ]
    if kind == "required_duration_hours":
        return [f"Asgari eğitim süresi: {_clean(block.get('value'))} saat"]
    if kind == "required_duration_minutes":
        return [f"Ders süresi karşılığı: {_clean(block.get('value'))} dakika"]
    if kind == "topic_count":
        return [f"İşe özgü konu sayısı: {_clean(block.get('value'))}"]
    if kind == "approved_foundation_questions":
        return [f"Onaylı temel soru sayısı: {_clean(block.get('count'))}"]
    if kind == "frozen_training_topic":
        return [f"Dondurulmuş konu {_clean(block.get('index'))}: {_clean(block.get('value'))}"]
    if kind == "technical_content_pending_renderer":
        return ["İçerik, bu eğitime ait dondurulmuş konu ve risk kayıtlarıyla sınırlandırılmıştır."]
    if kind == "technical_risk_tag":
        return [f"Teknik risk: {_risk_label(block.get('value'))}"]
    if kind == "special_risk":
        return [f"Özel risk: {_risk_label(block.get('value'))}"]
    if kind == "special_risks_empty":
        return ["Özel risk listesi boş; ek risk uydurulmamıştır."]
    if kind == "control_hierarchy":
        return [
            "Kontrol sırası: ortadan kaldırma veya azaltma",
            "Mühendislik ve toplu korunma tedbirleri",
            "Organizasyon, talimat ve gözetim",
            "Uygun KKD - son kontrol katmanı",
        ]
    if kind == "topic_control_mapping":
        return [f"Konu-kontrol eşleştirmesi: {_clean(block.get('topic'))}"]
    if kind == "ppe_mapping_pending_specialist_approval":
        risks = [_risk_label(item) for item in (block.get("risk_tags") or [])]
        text = ", ".join(risks[:5])
        return [
            "KKD seçimi işyeri risk değerlendirmesiyle doğrulanmalıdır.",
            f"İlişkili riskler: {text}" if text else "İlişkili risk kaydı bulunmuyor.",
        ]
    if kind == "workplace_emergency_placeholder":
        field = str(block.get("field") or "")
        label = {
            "assembly_point": "Toplanma yeri",
            "emergency_contacts": "Acil iletişim bilgileri",
        }.get(field, field.replace("_", " ").title())
        return [f"{label}: işyerine özel bilgi ve uzman onayı bekleniyor"]
    if kind == "exam_distribution":
        return [
            f"Temel soru: {_clean(block.get('foundation'))}",
            f"İşe özgü soru: {_clean(block.get('work_specific'))}",
        ]
    if kind == "exam_workflow_unchanged":
        return ["Mevcut sınav, puanlama ve sertifika iş akışı değiştirilmemiştir."]
    if kind == "topic_summary":
        return [f"Konu: {_clean(item)}" for item in (block.get("values") or [])]
    if kind == "risk_summary":
        return [f"Risk: {_risk_label(item)}" for item in (block.get("values") or [])]
    if kind in {"contract_version", "template_version", "catalog_version", "catalog_hash"}:
        label = kind.replace("_", " ").title()
        return [f"{label}: {_clean(block.get('value'), limit=100)}"]
    if kind:
        value = block.get("value")
        return [f"{kind.replace('_', ' ').title()}: {_clean(value)}"] if value not in (None, "") else []
    return []


def slide_bullets(slide: dict[str, Any], *, max_items: int = 6) -> list[str]:
    bullets: list[str] = []
    for block in slide.get("content_blocks") or []:
        if not isinstance(block, dict):
            continue
        for item in _block_bullets(block):
            item = _clean(item, limit=220)
            if item and item not in bullets:
                bullets.append(item)
    if len(bullets) <= max_items:
        return bullets
    kept = bullets[: max_items - 1]
    kept.append(f"Ayrıntılı manifestte {len(bullets) - len(kept)} ek kayıt bulunmaktadır.")
    return kept


def _add_textbox(slide, x, y, w, h, text: str, *, size: float, color: RGBColor, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font_name="Aptos"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def _add_bullet_box(slide, bullets: list[str], *, x=0.9, y=1.75, w=11.55, h=4.65):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    for index, text in enumerate(bullets or ["Bu slayt için görüntülenecek içerik bulunmuyor."]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {text}"
        paragraph.space_after = Pt(9)
        paragraph.line_spacing = 1.1
        for run in paragraph.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(20 if len(bullets) <= 4 else 17)
            run.font.color.rgb = _NAVY
    return box


def _pptx_slide(prs: Presentation, manifest: dict[str, Any], item: dict[str, Any], total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _WHITE

    position = int(item["position"])
    section = _clean(item.get("section_id"), limit=60)
    title = _clean(item.get("title"), limit=140)
    sources = _source_titles(manifest, list(item.get("source_refs") or []))

    if section == "cover":
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        shape.fill.solid(); shape.fill.fore_color.rgb = _TEAL_DARK
        shape.line.fill.background()
        _add_textbox(slide, 0.85, 0.72, 11.7, 0.55, "İSG SUITE OSGB", size=18, color=_MINT, bold=True)
        _add_textbox(slide, 0.85, 1.55, 11.45, 1.45, title, size=34, color=_WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
        bullets = slide_bullets(item, max_items=5)
        _add_textbox(slide, 0.9, 3.35, 11.1, 1.8, "\n".join(bullets), size=18, color=_WHITE)
        _add_textbox(slide, 0.9, 6.62, 11.4, 0.35, f"Sürüm: {manifest.get('template_version')} · Manifest: {str(manifest.get('content_hash'))[:12]}", size=10, color=_MINT)
    else:
        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.24))
        top.fill.solid(); top.fill.fore_color.rgb = _TEAL
        top.line.fill.background()
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(0.56), Inches(2.3), Inches(0.38))
        badge.fill.solid(); badge.fill.fore_color.rgb = _MINT
        badge.line.fill.background()
        tf = badge.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = section.replace("_", " ").upper(); p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = "Aptos"; run.font.size = Pt(10); run.font.bold = True; run.font.color.rgb = _TEAL_DARK
        _add_textbox(slide, 0.72, 1.02, 11.8, 0.63, title, size=27, color=_NAVY, bold=True)
        _add_bullet_box(slide, slide_bullets(item))
        if item.get("approval_required"):
            callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.68), Inches(0.55), Inches(2.55), Inches(0.42))
            callout.fill.solid(); callout.fill.fore_color.rgb = RGBColor(255, 242, 204)
            callout.line.color.rgb = RGBColor(230, 190, 90)
            tf = callout.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.text = "UZMAN ONAYI GEREKLİ"; p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = "Aptos"; run.font.size = Pt(9); run.font.bold = True; run.font.color.rgb = _AMBER

        footer = "Kaynak: " + (" · ".join(sources[:2]) if sources else "Dondurulmuş sunum manifesti")
        _add_textbox(slide, 0.72, 7.02, 10.9, 0.24, _clean(footer, limit=180), size=8.5, color=_SLATE)
        _add_textbox(slide, 12.0, 7.0, 0.65, 0.25, f"{position}/{total}", size=9, color=_SLATE, bold=True, align=PP_ALIGN.RIGHT)

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        f"Slayt {position}/{total}: {title}\n"
        f"Bölüm: {section}\n"
        f"Kaynaklar: {'; '.join(sources) if sources else 'Dondurulmuş manifest'}\n"
        "Bu slayt yalnız sürüm kaydındaki değişmez manifestten üretilmiştir."
    )
    return slide


def _normalize_zip(data: bytes) -> bytes:
    source = BytesIO(data)
    output = BytesIO()
    with zipfile.ZipFile(source, "r") as zin, zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9) as zout:
        for original in zin.infolist():
            info = zipfile.ZipInfo(original.filename, date_time=(1980, 1, 1, 0, 0, 0))
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = original.external_attr
            info.create_system = 0
            zout.writestr(info, zin.read(original.filename))
    return output.getvalue()


def render_pptx(manifest: dict[str, Any]) -> bytes:
    verify_manifest(manifest)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    props = prs.core_properties
    props.title = _clean((manifest.get("training") or {}).get("title") or "NACE Uyumlu İSG Eğitimi", limit=120)
    props.subject = "NACE uyumlu iş sağlığı ve güvenliği eğitim sunumu"
    props.author = "İSG Suite OSGB"
    props.last_modified_by = "İSG Suite OSGB"
    props.created = datetime(2000, 1, 1)
    props.modified = datetime(2000, 1, 1)
    props.comments = f"{RENDERER_VERSION}; manifest={manifest.get('content_hash')}"
    slides = list(manifest["slides"])
    for item in slides:
        _pptx_slide(prs, manifest, item, len(slides))
    buffer = BytesIO()
    prs.save(buffer)
    result = _normalize_zip(buffer.getvalue())
    check = Presentation(BytesIO(result))
    if len(check.slides) != len(slides):
        raise RuntimeError("PPTX slayt sayısı doğrulanamadı.")
    if not zipfile.is_zipfile(BytesIO(result)):
        raise RuntimeError("PPTX paket doğrulaması başarısız.")
    return result


_PDF_FONT = "Helvetica"
_PDF_FONT_B = "Helvetica-Bold"


def _register_pdf_fonts() -> None:
    global _PDF_FONT, _PDF_FONT_B
    candidates = [
        (
            Path(__file__).resolve().parents[1] / "assets" / "fonts" / "DejaVuSans.ttf",
            Path(__file__).resolve().parents[1] / "assets" / "fonts" / "DejaVuSans-Bold.ttf",
        ),
        (Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"), Path("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf")),
        (Path("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"), Path("/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf")),
    ]
    for regular, bold in candidates:
        if regular.exists():
            pdfmetrics.registerFont(TTFont("PresentationSans", str(regular)))
            pdfmetrics.registerFont(TTFont("PresentationSans-Bold", str(bold if bold.exists() else regular)))
            _PDF_FONT, _PDF_FONT_B = "PresentationSans", "PresentationSans-Bold"
            return
    raise RuntimeError("Sunum PDF için Unicode font bulunamadı.")


def _pdf_wrap(text: str, width: float, font: str, size: float, *, max_lines: int = 3) -> list[str]:
    words = _clean(text, limit=500).split()
    lines: list[str] = []
    current = ""
    for word in words:
        trial = f"{current} {word}".strip()
        if pdfmetrics.stringWidth(trial, font, size) <= width:
            current = trial
        else:
            if current:
                lines.append(current)
            current = word
            if len(lines) >= max_lines:
                break
    if current and len(lines) < max_lines:
        lines.append(current)
    if len(lines) == max_lines and len(" ".join(lines)) < len(" ".join(words)):
        last = lines[-1]
        while last and pdfmetrics.stringWidth(last + "…", font, size) > width:
            last = last[:-1]
        lines[-1] = last.rstrip() + "…"
    return lines or [""]


def render_pdf(manifest: dict[str, Any]) -> bytes:
    verify_manifest(manifest)
    _register_pdf_fonts()
    page_size = (SLIDE_WIDTH_IN * inch, SLIDE_HEIGHT_IN * inch)
    buffer = BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=page_size, pageCompression=1, invariant=1)
    pdf.setTitle(_clean((manifest.get("training") or {}).get("title") or "NACE Uyumlu İSG Eğitimi", limit=120))
    pdf.setAuthor("İSG Suite OSGB")
    pdf.setCreator(RENDERER_VERSION)
    width, height = page_size
    slides = list(manifest["slides"])
    registry = _source_map(manifest)

    for item in slides:
        position = int(item["position"])
        section = _clean(item.get("section_id"), limit=60)
        title = _clean(item.get("title"), limit=140)
        if section == "cover":
            pdf.setFillColorRGB(6 / 255, 84 / 255, 82 / 255)
            pdf.rect(0, 0, width, height, fill=1, stroke=0)
            pdf.setFillColorRGB(224 / 255, 245 / 255, 242 / 255)
            pdf.setFont(_PDF_FONT_B, 17)
            pdf.drawString(0.85 * inch, height - 0.95 * inch, "İSG SUITE OSGB")
            pdf.setFillColorRGB(1, 1, 1)
            y = height - 1.7 * inch
            for line in _pdf_wrap(title, 11.2 * inch, _PDF_FONT_B, 30, max_lines=3):
                pdf.setFont(_PDF_FONT_B, 30); pdf.drawString(0.85 * inch, y, line); y -= 0.48 * inch
            pdf.setFont(_PDF_FONT, 16)
            y = height - 3.55 * inch
            for bullet in slide_bullets(item, max_items=5):
                for line_index, line in enumerate(_pdf_wrap(bullet, 10.8 * inch, _PDF_FONT, 16, max_lines=2)):
                    prefix = "• " if line_index == 0 else "  "
                    pdf.drawString(0.95 * inch, y, prefix + line)
                    y -= 0.32 * inch
                y -= 0.07 * inch
            pdf.setFillColorRGB(224 / 255, 245 / 255, 242 / 255)
            pdf.setFont(_PDF_FONT, 9)
            pdf.drawString(0.85 * inch, 0.48 * inch, f"Sürüm: {manifest.get('template_version')} · Manifest: {str(manifest.get('content_hash'))[:12]}")
        else:
            pdf.setFillColorRGB(1, 1, 1); pdf.rect(0, 0, width, height, fill=1, stroke=0)
            pdf.setFillColorRGB(8 / 255, 127 / 255, 120 / 255); pdf.rect(0, height - 0.24 * inch, width, 0.24 * inch, fill=1, stroke=0)
            pdf.setFillColorRGB(6 / 255, 84 / 255, 82 / 255); pdf.setFont(_PDF_FONT_B, 9)
            pdf.drawString(0.72 * inch, height - 0.77 * inch, section.replace("_", " ").upper())
            pdf.setFillColorRGB(17 / 255, 38 / 255, 53 / 255)
            y = height - 1.25 * inch
            for line in _pdf_wrap(title, 11.6 * inch, _PDF_FONT_B, 24, max_lines=2):
                pdf.setFont(_PDF_FONT_B, 24); pdf.drawString(0.72 * inch, y, line); y -= 0.36 * inch
            y = min(y - 0.18 * inch, height - 1.85 * inch)
            bullets = slide_bullets(item)
            size = 17 if len(bullets) <= 4 else 14.5
            for bullet in bullets or ["Bu slayt için görüntülenecek içerik bulunmuyor."]:
                for line_index, line in enumerate(_pdf_wrap(bullet, 11.1 * inch, _PDF_FONT, size, max_lines=2)):
                    pdf.setFont(_PDF_FONT, size)
                    pdf.drawString(0.9 * inch, y, ("• " if line_index == 0 else "  ") + line)
                    y -= (0.31 if size >= 16 else 0.27) * inch
                y -= 0.07 * inch
            if item.get("approval_required"):
                pdf.setFillColorRGB(180 / 255, 100 / 255, 0)
                pdf.setFont(_PDF_FONT_B, 9)
                pdf.drawRightString(width - 0.72 * inch, height - 0.75 * inch, "UZMAN ONAYI GEREKLİ")
            refs = list(item.get("source_refs") or [])
            source_titles = []
            for ref in refs:
                key = str(ref)
                source_titles.append(registry.get(key, key.split(":", 1)[0].replace("_", " ")))
            footer = "Kaynak: " + (" · ".join(source_titles[:2]) if source_titles else "Dondurulmuş sunum manifesti")
            pdf.setFillColorRGB(82 / 255, 101 / 255, 120 / 255); pdf.setFont(_PDF_FONT, 7.5)
            pdf.drawString(0.72 * inch, 0.25 * inch, _clean(footer, limit=180))
            pdf.setFont(_PDF_FONT_B, 8.5); pdf.drawRightString(width - 0.72 * inch, 0.25 * inch, f"{position}/{len(slides)}")
        pdf.showPage()
    pdf.save()
    result = buffer.getvalue()
    reader = PdfReader(BytesIO(result))
    if len(reader.pages) != len(slides):
        raise RuntimeError("PDF sayfa sayısı doğrulanamadı.")
    return result


def render_presentation(manifest: dict[str, Any]) -> RenderedPresentation:
    pptx_bytes = render_pptx(manifest)
    pdf_bytes = render_pdf(manifest)
    return RenderedPresentation(
        pptx_bytes=pptx_bytes,
        pdf_bytes=pdf_bytes,
        slide_count=int(manifest["slide_count"]),
    )
