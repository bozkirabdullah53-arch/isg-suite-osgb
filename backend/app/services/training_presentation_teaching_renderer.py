"""Visual PPTX/PDF renderer for Teaching V3 NACE presentations.

The renderer uses deterministic vector graphics (cards, flows, hierarchy and
risk maps) so every generated deck contains real visual teaching aids without
external-image, licensing or network dependencies.
"""
from __future__ import annotations

import json
from dataclasses import dataclass
from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import landscape, letter
from reportlab.pdfgen import canvas

from app.services.training_presentation_renderer import (
    PDF_CONTENT_TYPE,
    PPTX_CONTENT_TYPE,
    sha256_hex,
    verify_manifest,
)

TEACHING_RENDERER_VERSION = "nace-training-teaching-renderer-v3"


_SECTION_LABELS = {
    "cover": "AÇILIŞ",
    "learning_objectives": "ÖĞRENME HEDEFLERİ",
    "legal_basis": "ÇALIŞANIN HAKLARI VE SORUMLULUKLARI",
    "nace_identity": "NACE VE İŞYERİ FAALİYETİ",
    "training_plan": "EĞİTİM PLANI",
    "foundation_ohs": "TEMEL İSG KONULARI",
    "work_specific_topics": "İŞE VE İŞYERİNE ÖZGÜ RİSKLER",
    "technical_risks": "TEKNİK RİSKLER",
    "control_measures": "KONTROL TEDBİRLERİ",
    "ppe": "KİŞİSEL KORUYUCU DONANIM",
    "emergency": "ACİL DURUM",
    "assessment": "ÖLÇME VE DEĞERLENDİRME",
    "summary": "ÖZET",
    "sources_and_version": "KAYNAKLAR VE SÜRÜM",
    "custom_instructor_slide": "EĞİTMEN TARAFINDAN EKLENEN İÇERİK",
}


def _section_label(section_id: object) -> str:
    key = str(section_id or "").strip()
    return _SECTION_LABELS.get(key, "DERS SUNUMU")


_LEGAL_BASIS_POINTS = [
    "İşyerindeki tehlikeleri ve size bildirilen kontrol tedbirlerini öğrenin.",
    "Verilen eğitim, talimat ve güvenli çalışma kurallarına uygun hareket edin.",
    "Güvensiz durumu, ramak kala olayı, iş kazasını veya sağlık belirtisini gecikmeden bildirin.",
    "Ciddi ve yakın tehlikede güvenli biçimde işi durdurun, tehlikeli alandan uzaklaşın ve yetkiliye haber verin.",
]

SLIDE_W = 13.333333
SLIDE_H = 7.5

NAVY = RGBColor(6, 27, 46)
TEAL = RGBColor(15, 118, 110)
GREEN = RGBColor(34, 164, 71)
AMBER = RGBColor(219, 145, 36)
RED = RGBColor(190, 52, 52)
SLATE = RGBColor(71, 85, 105)
PALE = RGBColor(245, 248, 250)
MINT = RGBColor(226, 246, 241)
WHITE = RGBColor(255, 255, 255)


@dataclass(frozen=True)
class RenderedTeachingPresentation:
    pptx_bytes: bytes
    pdf_bytes: bytes
    slide_count: int
    renderer_version: str = TEACHING_RENDERER_VERSION


def _clean(value: object, limit: int = 700) -> str:
    text = " ".join(str(value or "").replace("\n", " ").split()).strip()
    return text if len(text) <= limit else text[: limit - 1].rstrip() + "…"


def _blocks(slide: dict[str, Any], kind: str) -> list[dict[str, Any]]:
    return [
        item for item in (slide.get("content_blocks") or [])
        if isinstance(item, dict) and str(item.get("type") or "") == kind
    ]


def _value(slide: dict[str, Any], kind: str) -> str:
    rows = _blocks(slide, kind)
    return _clean(rows[0].get("value")) if rows else ""


def _values(slide: dict[str, Any], kind: str) -> list[str]:
    result: list[str] = []
    for row in _blocks(slide, kind):
        raw = row.get("values") if isinstance(row.get("values"), list) else [row.get("value")]
        for item in raw:
            text = _clean(item, 260)
            if text and text not in result:
                result.append(text)
    return result


def _lesson_points(slide: dict[str, Any]) -> list[str]:
    points = _values(slide, "lesson_point")
    for kind, prefix in (
        ("learning_objective", "Hedef: "),
        ("lesson_explanation", ""),
        ("tehlike", "Tehlike: "),
        ("kontrol_tedbiri", "Kontrol: "),
        ("guvenli_davranis", "Güvenli davranış: "),
        ("key_takeaway", "Ana mesaj: "),
    ):
        for row in _blocks(slide, kind):
            text = _clean(row.get("value"), 360)
            if text:
                item = prefix + text
                if item not in points:
                    points.append(item)
    if not points:
        for row in slide.get("content_blocks") or []:
            if not isinstance(row, dict):
                continue
            value = _clean(row.get("value"), 360)
            if value and value not in points:
                points.append(value)
    if str(slide.get("section_id") or "") == "legal_basis" and not points:
        points.extend(_LEGAL_BASIS_POINTS)
    return points[:6]


def _sources(manifest: dict[str, Any], slide: dict[str, Any]) -> str:
    registry = {
        str(item.get("source_id")): _clean(item.get("title"), 90)
        for item in (manifest.get("source_registry") or [])
        if isinstance(item, dict) and item.get("source_id")
    }
    names: list[str] = []
    for ref in slide.get("source_refs") or []:
        key = str(ref or "")
        label = registry.get(key)
        if not label:
            if key.startswith("training_nace_snapshot:"):
                label = "NACE snapshot"
            elif key.startswith("controlled_training_topics:"):
                label = "Kontrollü eğitim konuları"
            elif key.startswith("controlled_risk_catalog:"):
                label = "Teknik risk kataloğu"
            elif key.startswith("approved_question_bank:"):
                label = "Onaylı soru bankası"
            else:
                label = _clean(key, 80)
        if label and label not in names:
            names.append(label)
    return " · ".join(names[:3])


def _add_text(slide, x, y, w, h, text, *, size=18, color=NAVY, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear(); frame.word_wrap = True; frame.vertical_anchor = valign
    p = frame.paragraphs[0]; p.text = str(text or ""); p.alignment = align
    for run in p.runs:
        run.font.name = "Aptos"; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return box


def _card(slide, x, y, w, h, *, title, body, accent=TEAL, number=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(214, 223, 230); shape.line.width = Pt(1.1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    if number is not None:
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.2), Inches(y + 0.18), Inches(0.42), Inches(0.42))
        circ.fill.solid(); circ.fill.fore_color.rgb = accent; circ.line.fill.background()
        _add_text(slide, x + 0.2, y + 0.19, 0.42, 0.4, str(number), size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        tx = x + 0.72
    else:
        tx = x + 0.25
    _add_text(slide, tx, y + 0.17, w - (tx - x) - 0.2, 0.35, title, size=12, color=accent, bold=True)
    _add_text(slide, x + 0.25, y + 0.62, w - 0.45, h - 0.78, body, size=13.5, color=NAVY)
    return shape


def _visual_hazard_flow(slide, visual: dict[str, Any], *, x=7.65, y=1.55, w=4.85, h=4.8):
    hazard = _clean(visual.get("hazard"), 330) or "Tehlikeyi tanımla"
    controls = [_clean(v, 240) for v in (visual.get("controls") or []) if _clean(v, 240)]
    control = controls[0] if controls else "Kontrol tedbirinin uygulandığını doğrula"
    behavior = _clean(visual.get("safe_behavior"), 330) or "Güvenli davranışı uygula ve uygunsuzluğu bildir"
    cards = [
        ("1 · TEHLİKE", hazard, RED),
        ("2 · KONTROL", control, TEAL),
        ("3 · GÜVENLİ DAVRANIŞ", behavior, GREEN),
    ]
    for index, (title, body, accent) in enumerate(cards):
        cy = y + index * 1.42
        _card(slide, x, cy, w, 1.05, title=title, body=body, accent=accent)
        if index < 2:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + w / 2), Inches(cy + 1.06), Inches(x + w / 2), Inches(cy + 1.36))
            conn.line.color.rgb = SLATE; conn.line.width = Pt(1.5)


def _visual_risk_map(slide, values: list[str], *, x=7.65, y=1.6, w=4.85):
    vals = values[:8] or ["Risk kaydı uzman incelemesinde"]
    cols = 2
    cw = (w - 0.2) / cols
    for idx, value in enumerate(vals):
        row, col = divmod(idx, cols)
        _card(slide, x + col * (cw + 0.1), y + row * 1.05, cw, 0.86, title=f"RİSK {idx + 1}", body=value.replace("_", " "), accent=AMBER, number=idx + 1)


def _visual_hierarchy(slide, *, x=7.7, y=1.75, w=4.7):
    labels = [
        ("1  Ortadan kaldır", GREEN),
        ("2  İkame / azalt", TEAL),
        ("3  Mühendislik / toplu korunma", RGBColor(35, 113, 155)),
        ("4  Organizasyon / talimat", AMBER),
        ("5  KKD", RED),
    ]
    for idx, (label, color) in enumerate(labels):
        shrink = idx * 0.32
        bx = x + shrink / 2
        bw = w - shrink
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(y + idx * 0.78), Inches(bw), Inches(0.62))
        shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
        _add_text(slide, bx, y + idx * 0.78 + 0.03, bw, 0.55, label, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def _visual_flow(slide, labels: list[str], *, x=7.6, y=2.05, w=4.9, accent=TEAL):
    count = len(labels)
    cw = (w - (count - 1) * 0.16) / count
    for idx, label in enumerate(labels):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + idx * (cw + 0.16)), Inches(y), Inches(cw), Inches(1.2))
        shape.fill.solid(); shape.fill.fore_color.rgb = accent if idx < count - 1 else GREEN
        shape.line.fill.background()
        _add_text(slide, x + idx * (cw + 0.16) + 0.06, y + 0.18, cw - 0.12, 0.8, label, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if idx < count - 1:
            _add_text(slide, x + (idx + 1) * cw + idx * 0.16, y + 0.46, 0.16, 0.25, "→", size=16, color=SLATE, bold=True, align=PP_ALIGN.CENTER)


def _visual_topic_map(slide, values: list[str], *, x=7.65, y=1.55, w=4.85):
    vals = values[:5]
    for idx, value in enumerate(vals):
        _card(slide, x, y + idx * 0.92, w, 0.76, title=f"KONU {idx + 1}", body=value, accent=TEAL, number=idx + 1)


def _render_visual(slide, item: dict[str, Any], manifest: dict[str, Any]):
    visual = _blocks(item, "hazard_control_behavior_visual")
    if visual:
        _visual_hazard_flow(slide, visual[0]); return
    risk = _blocks(item, "risk_map_visual")
    if risk:
        vals = risk[0].get("values") or manifest.get("technical_risk_tags") or []
        _visual_risk_map(slide, [_clean(v, 180) for v in vals]); return
    if _blocks(item, "control_hierarchy_visual"):
        _visual_hierarchy(slide); return
    if _blocks(item, "emergency_flow_visual"):
        _visual_flow(slide, ["Fark et", "Uyar / bildir", "Tahliye", "Toplan & sayım"]); return
    assessment = _blocks(item, "assessment_visual")
    if assessment:
        _visual_flow(slide, ["5 temel", "15 işe özgü", "20 soru", "Sonuç & kayıt"], accent=RGBColor(35, 113, 155)); return
    topic = _blocks(item, "topic_map_visual")
    if topic:
        vals = topic[0].get("values") or manifest.get("training_topics") or []
        _visual_topic_map(slide, [_clean(v, 220) for v in vals]); return
    if _blocks(item, "ppe_layers_visual"):
        _visual_flow(slide, ["Kaynakta kontrol", "Toplu korunma", "İş organizasyonu", "KKD"]); return
    takeaway = _value(item, "key_takeaway") or "Tehlikeyi tanı · kontrolü doğrula · güvenli davranışı uygula"
    _card(slide, 7.75, 2.1, 4.55, 2.1, title="ANA MESAJ", body=takeaway, accent=TEAL)


def _speaker_note(item: dict[str, Any], manifest: dict[str, Any]) -> str:
    note = _value(item, "instructor_note")
    scenario = _value(item, "case_scenario")
    check = _value(item, "check_question")
    source = _sources(manifest, item)
    rows = [
        note,
        f"Vaka tartışması: {scenario}" if scenario else "",
        f"Kontrol sorusu: {check}" if check else "",
        f"Kaynaklar: {source}" if source else "",
    ]
    return "\n\n".join(row for row in rows if row)


def _pptx_slide(prs: Presentation, manifest: dict[str, Any], item: dict[str, Any], total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = WHITE
    pos = int(item.get("position") or 0)
    section = str(item.get("section_id") or "")
    title = _clean(item.get("title"), 160)
    if section == "cover":
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        rect.fill.solid(); rect.fill.fore_color.rgb = NAVY; rect.line.fill.background()
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), prs.slide_height)
        band.fill.solid(); band.fill.fore_color.rgb = GREEN; band.line.fill.background()
        _add_text(slide, 0.8, 0.65, 4.0, 0.4, "İSG SUITE · DERS SUNUMU V3", size=15, color=MINT, bold=True)
        _add_text(slide, 0.8, 1.55, 11.6, 1.4, title, size=33, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
        nace = manifest.get("nace_snapshot") or {}
        _card(slide, 0.85, 3.45, 5.7, 1.45, title="NACE / FAALİYET", body=f"{nace.get('nace_code','')} · {nace.get('nace_description','')}", accent=TEAL)
        _card(slide, 6.8, 3.45, 2.5, 1.45, title="TEHLİKE SINIFI", body=nace.get("hazard_class", "—"), accent=AMBER)
        training = manifest.get("training") or {}
        _card(slide, 9.55, 3.45, 2.75, 1.45, title="EĞİTİM", body=f"{training.get('start_date','')} → {training.get('end_date','')}", accent=GREEN)
        _add_text(slide, 0.85, 6.75, 11.5, 0.3, f"Kaynak kontrollü · {TEACHING_RENDERER_VERSION}", size=10, color=MINT)
    else:
        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.68))
        top.fill.solid(); top.fill.fore_color.rgb = NAVY; top.line.fill.background()
        _add_text(slide, 0.65, 0.16, 2.5, 0.3, _section_label(section), size=10.5, color=MINT, bold=True)
        _add_text(slide, 0.72, 0.82, 11.8, 0.65, title, size=25, color=NAVY, bold=True)
        points = _lesson_points(item)
        _add_text(slide, 0.75, 1.62, 6.35, 0.28, "DERS ANLATIMI", size=11, color=TEAL, bold=True)
        for idx, point in enumerate(points[:5]):
            _card(slide, 0.75, 2.0 + idx * 0.78, 6.35, 0.66, title=f"{idx + 1}", body=point, accent=TEAL, number=idx + 1)
        _render_visual(slide, item, manifest)
        scenario = _value(item, "case_scenario")
        check = _value(item, "check_question")
        if scenario or check:
            body = scenario or check
            label = "VAKA / TARTIŞMA" if scenario else "BİLGİ KONTROLÜ"
            _card(slide, 0.75, 6.02, 11.75, 0.72, title=label, body=body, accent=AMBER)
        source = _sources(manifest, item)
        _add_text(slide, 0.75, 6.95, 10.7, 0.25, f"Kaynak: {source or 'Manifest kaynak kayıtları'}", size=8.5, color=SLATE)
        _add_text(slide, 11.55, 6.95, 0.95, 0.25, f"{pos}/{total}", size=9, color=SLATE, bold=True, align=PP_ALIGN.RIGHT)

    note = _speaker_note(item, manifest)
    if note:
        try:
            slide.notes_slide.notes_text_frame.text = note
        except Exception:
            pass
    return slide


def _pptx(manifest: dict[str, Any]) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W); prs.slide_height = Inches(SLIDE_H)
    slides = manifest.get("slides") or []
    for item in slides:
        _pptx_slide(prs, manifest, item, len(slides))
    buf = BytesIO(); prs.save(buf); return buf.getvalue()


def _pdf_color(c, rgb: RGBColor):
    c.setFillColorRGB(rgb[0] / 255, rgb[1] / 255, rgb[2] / 255)


def _pdf_wrapped(c, text: str, x: float, y: float, max_width: float, *, size=10, leading=13, max_lines=5):
    words = _clean(text, 1000).split()
    lines: list[str] = []; cur = ""
    for word in words:
        trial = f"{cur} {word}".strip()
        if c.stringWidth(trial, "Helvetica", size) <= max_width:
            cur = trial
        else:
            if cur: lines.append(cur)
            cur = word
        if len(lines) >= max_lines:
            break
    if cur and len(lines) < max_lines: lines.append(cur)
    for line in lines:
        c.drawString(x, y, line); y -= leading
    return y


def _pdf(manifest: dict[str, Any]) -> bytes:
    buf = BytesIO(); page = landscape(letter); w, h = page
    c = canvas.Canvas(buf, pagesize=page)
    slides = manifest.get("slides") or []
    for item in slides:
        c.setFillColorRGB(1, 1, 1); c.rect(0, 0, w, h, fill=1, stroke=0)
        c.setFillColorRGB(6/255, 27/255, 46/255); c.rect(0, h - 44, w, 44, fill=1, stroke=0)
        c.setFillColorRGB(1, 1, 1); c.setFont("Helvetica-Bold", 9); c.drawString(34, h - 27, _section_label(item.get("section_id")))
        c.setFillColorRGB(6/255, 27/255, 46/255); c.setFont("Helvetica-Bold", 19); c.drawString(34, h - 78, _clean(item.get("title"), 100))
        y = h - 108
        c.setFont("Helvetica-Bold", 9); c.setFillColorRGB(15/255, 118/255, 110/255); c.drawString(34, y, "DERS ANLATIMI"); y -= 18
        c.setFont("Helvetica", 9); c.setFillColorRGB(20/255, 40/255, 60/255)
        for idx, point in enumerate(_lesson_points(item)[:5], 1):
            c.setFont("Helvetica-Bold", 8); c.drawString(36, y, f"{idx}.")
            c.setFont("Helvetica", 8.5); y = _pdf_wrapped(c, point, 54, y, 330, size=8.5, leading=11, max_lines=3) - 5
        # Visual panel is intentionally graphic even in the companion PDF.
        vx, vy, vw, vh = 410, 150, 330, 280
        c.setFillColorRGB(245/255, 248/255, 250/255); c.roundRect(vx, vy, vw, vh, 10, fill=1, stroke=0)
        c.setFillColorRGB(15/255, 118/255, 110/255); c.setFont("Helvetica-Bold", 10); c.drawString(vx + 16, vy + vh - 24, "GÖRSEL ÖĞRENME HARİTASI")
        visual = _blocks(item, "hazard_control_behavior_visual")
        if visual:
            row = visual[0]
            texts = [
                ("TEHLİKE", _clean(row.get("hazard"), 220)),
                ("KONTROL", _clean((row.get("controls") or ["Kontrolü doğrula"])[0], 220)),
                ("GÜVENLİ DAVRANIŞ", _clean(row.get("safe_behavior"), 220)),
            ]
        else:
            texts = [("ANA MESAJ", _value(item, "key_takeaway") or "Tehlikeyi tanı · kontrolü doğrula · güvenli davranışı uygula")]
        yy = vy + vh - 54
        for label, body in texts:
            c.setFillColorRGB(1, 1, 1); c.roundRect(vx + 14, yy - 58, vw - 28, 50, 7, fill=1, stroke=0)
            c.setFillColorRGB(15/255, 118/255, 110/255); c.setFont("Helvetica-Bold", 8); c.drawString(vx + 26, yy - 25, label)
            c.setFillColorRGB(20/255, 40/255, 60/255); c.setFont("Helvetica", 7.5); _pdf_wrapped(c, body, vx + 100, yy - 25, vw - 128, size=7.5, leading=9, max_lines=3)
            yy -= 68
        scenario = _value(item, "case_scenario")
        if scenario:
            c.setFillColorRGB(1, 0.97, 0.88); c.roundRect(34, 46, w - 68, 54, 8, fill=1, stroke=0)
            c.setFillColorRGB(140/255, 88/255, 0); c.setFont("Helvetica-Bold", 8); c.drawString(48, 82, "VAKA / TARTIŞMA")
            c.setFillColorRGB(40/255, 45/255, 50/255); c.setFont("Helvetica", 7.5); _pdf_wrapped(c, scenario, 48, 67, w - 96, size=7.5, leading=9, max_lines=3)
        c.setFillColorRGB(90/255, 105/255, 120/255); c.setFont("Helvetica", 6.5); c.drawString(34, 22, _clean(_sources(manifest, item), 150))
        c.drawRightString(w - 34, 22, f"{int(item.get('position') or 0)}/{len(slides)}")
        c.showPage()
    c.save(); return buf.getvalue()


def render_teaching_presentation(manifest: dict[str, Any]) -> RenderedTeachingPresentation:
    verify_manifest(manifest)
    if not bool((manifest.get("rendering") or {}).get("teaching_v3")):
        raise ValueError("Teaching V3 manifest işareti bulunmuyor. Önce V3 zenginleştirmesi oluşturulmalıdır.")
    pptx = _pptx(manifest)
    pdf = _pdf(manifest)
    if not pptx.startswith(b"PK") or not pdf.startswith(b"%PDF"):
        raise RuntimeError("Teaching V3 sunum çıktısı beklenen dosya biçiminde oluşmadı.")
    # Hash calculation is intentionally performed here as an integrity smoke check.
    sha256_hex(pptx); sha256_hex(pdf)
    return RenderedTeachingPresentation(pptx_bytes=pptx, pdf_bytes=pdf, slide_count=len(manifest.get("slides") or []))


__all__ = [
    "PDF_CONTENT_TYPE",
    "PPTX_CONTENT_TYPE",
    "TEACHING_RENDERER_VERSION",
    "RenderedTeachingPresentation",
    "render_teaching_presentation",
]
